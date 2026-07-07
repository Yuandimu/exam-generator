"""
智能出题系统 - 服务器部署版 v3（增强诊断）
支持 PPTX / PDF / DOCX / TXT / MD 文件上传与内容提取
适用于公司内部服务器部署
"""
import http.server
import json
import os
import re
import sys
import socket
import threading
import time
from io import BytesIO

# ── 配置 ──
PORT = 8080               # 服务器端口，可根据需要修改
HOST = '0.0.0.0'          # 监听所有网络接口（允许局域网访问）
DIRECTORY = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(DIRECTORY, "uploads")

os.makedirs(UPLOAD_DIR, exist_ok=True)


def extract_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        ct = cell.text.strip()
                        if ct:
                            row_text.append(ct)
                    if row_text:
                        slide_lines.append(" | ".join(row_text))
        if slide_lines:
            slides_text.append(f"--- Page {idx} ---\n" + "\n".join(slide_lines))
    return "\n\n".join(slides_text)


def extract_pdf(filepath):
    from PyPDF2 import PdfReader
    reader = PdfReader(filepath)
    pages_text = []
    for idx, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages_text.append(f"--- Page {idx} ---\n" + text.strip())
    return "\n\n".join(pages_text)


def extract_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                paragraphs.append(" | ".join(row_text))
    return "\n".join(paragraphs)


def extract_text(filepath, ext):
    if ext in ('.txt', '.md'):
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    elif ext == '.pptx':
        return extract_pptx(filepath)
    elif ext == '.pdf':
        return extract_pdf(filepath)
    elif ext == '.docx':
        return extract_docx(filepath)
    else:
        return None


def get_local_ip():
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        ip = s.getsockname()[0]
        s.close()
        return ip
    except Exception:
        return "unknown"


def open_firewall_windows(port):
    """On Windows: attempt to add firewall rule for the port"""
    if sys.platform != 'win32':
        return False
    try:
        import subprocess
        result = subprocess.run(
            ['netsh', 'advfirewall', 'firewall', 'add', 'rule',
             f'name=ExamGenerator Port {port}',
             'dir=in', 'action=allow',
             f'tcp', f'localport={str(port)}'],
            capture_output=True, timeout=10
        )
        if result.returncode == 0:
            print(f"  [OK] Firewall rule added for port {port}")
            return True
        else:
            err = result.stderr.decode('utf-8', errors='replace').strip()
            print(f"  [WARN] Firewall rule may already exist or need admin: {err}")
            # Rule might already exist - not a fatal error
            return True
    except Exception as e:
        print(f"  [WARN] Could not configure firewall: {e}")
        print(f"  [INFO] Please manually allow port {port} in Windows Firewall")
        return False


class ExamServer(http.server.SimpleHTTPRequestHandler):

    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=DIRECTORY, **kwargs)

    def end_headers(self):
        self.send_header('Cache-Control', 'no-store, no-cache, must-revalidate, max-age=0')
        self.send_header('Pragma', 'no-cache')
        self.send_header('Expires', '0')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type, Authorization')
        super().end_headers()

    def log_message(self, format, *args):
        if args and ('POST' in str(args) or 'error' in str(args).lower()):
            print(f"  [API] {args[0]}")

    def do_GET(self):
        # Diagnostic: log access attempts
        client = self.client_address[0]
        path = self.path
        print(f"  [GET] {client} -> {path}")
        super().do_GET()

    def do_POST(self):
        if self.path == '/api/upload':
            try:
                content_type = self.headers.get('Content-Type', '')
                if 'multipart/form-data' not in content_type:
                    self.send_error(400, "Need multipart/form-data")
                    return

                boundary = content_type.split('boundary=')[1].strip()
                if boundary.startswith('"') and boundary.endswith('"'):
                    boundary = boundary[1:-1]

                content_length = int(self.headers.get('Content-Length', 0))
                body = self.rfile.read(content_length)

                filename, file_data = parse_multipart(body, boundary)

                if not filename or not file_data:
                    self.send_json(400, {"error": "No file detected"})
                    return

                ext = os.path.splitext(filename)[1].lower()
                tmp_path = os.path.join(UPLOAD_DIR, filename)
                with open(tmp_path, 'wb') as f:
                    f.write(file_data)

                try:
                    content = extract_text(tmp_path, ext)
                except Exception as e:
                    self.send_json(500, {"error": f"Parse failed: {str(e)}"})
                    return
                finally:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)

                if content is None:
                    self.send_json(400, {"error": f"Unsupported format: {ext}"})
                    return

                if not content.strip():
                    self.send_json(400, {"error": "File contains no text (image-only?)"})

                print(f"  [API] Parsed {filename}: {len(content)} chars")
                self.send_json(200, {
                    "filename": filename,
                    "content": content,
                    "wordCount": len(content),
                    "ext": ext
                })

            except Exception as e:
                print(f"  [API] Upload error: {e}")
                import traceback
                traceback.print_exc()
                self.send_json(500, {"error": f"Server error: {str(e)}"})
        else:
            self.send_error(404, "Not found")

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def send_json(self, status, data):
        self.send_response(status)
        self.send_header('Content-Type', 'application/json; charset=utf-8')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))


def parse_multipart(body, boundary):
    boundary_bytes = boundary.encode('utf-8')
    parts = body.split(b'--' + boundary_bytes)

    for part in parts:
        if b'Content-Disposition' not in part:
            continue

        header_end = part.find(b'\r\n\r\n')
        if header_end == -1:
            continue

        header = part[:header_end].decode('utf-8', errors='replace')
        file_data = part[header_end + 4:]

        file_data = file_data.rstrip(b'\r\n')
        if file_data.endswith(b'--'):
            file_data = file_data[:-2].rstrip(b'\r\n')

        filename_match = re.search(r'filename="(.+?)"', header)
        if filename_match:
            return filename_match.group(1), file_data

    return None, None


def find_free_port(preferred_port):
    try:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
            s.bind((HOST, preferred_port))
            return preferred_port
    except OSError:
        print(f"  [WARN] Port {preferred_port} occupied, finding free port...")
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.setsockopt(socket.SOL_SOCKET, socket.SO_REUSEADDR, 1)
            s.bind((HOST, 0))
            return s.getsockname()[1]


if __name__ == "__main__":
    # ── Pre-flight checks ──
    print()
    print("========================================")
    print("  Exam Generator Server - Starting...")
    print("========================================")
    print()

    # Check index.html exists
    index_path = os.path.join(DIRECTORY, 'index.html')
    if not os.path.exists(index_path):
        print(f"[FATAL] index.html NOT FOUND at: {index_path}")
        print("Please ensure this script is in the same directory as index.html")
        input("Press Enter to exit...")
        sys.exit(1)
    print(f"[OK] Found index.html ({os.path.getsize(index_path)} bytes)")
    print(f"[OK] Working directory: {DIRECTORY}")

    # Find available port
    port = find_free_port(PORT)
    local_ip = get_local_ip()

    # Try to configure firewall on Windows
    if sys.platform == 'win32':
        print(f"\n[*] Attempting to configure Windows Firewall for port {port}...")
        open_firewall_windows(port)

    print()
    print("=" * 50)
    print("  SERVER IS RUNNING!")
    print("=" * 50)
    print(f"  Local:     http://localhost:{port}")
    print(f"  LAN:       http://{local_ip}:{port}")
    print(f"  Host:      {HOST}:{port}")
    print("-" * 50)
    print("  Press Ctrl+C to stop")
    print("=" * 50)
    print()
    print("[INFO] If you cannot open the page from another PC:")
    print(f"  1. Check Windows Firewall allows port {port}")
    print(f"  2. Temporarily disable antivirus/firewall to test")
    print(f"  3. Ensure both PCs are on the same network")
    print()

    try:
        server = http.server.HTTPServer((HOST, port), ExamServer)
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nServer stopped.")
        sys.exit(0)
    except Exception as e:
        print(f"\n[FATAL] Server crashed: {e}")
        import traceback
        traceback.print_exc()
        input("Press Enter to exit...")
        sys.exit(1)
