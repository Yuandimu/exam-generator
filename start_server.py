"""
智能出题系统 - 本地服务器（含文件解析API）
支持 PPTX / PDF / DOCX / TXT / MD 文件上传与内容提取
"""
import http.server
import json
import os
import re
import sys
import socket
import tempfile
import threading
import time
import webbrowser
from io import BytesIO

PORT = 3399
DIRECTORY = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(DIRECTORY, "uploads")

# 确保上传目录存在
os.makedirs(UPLOAD_DIR, exist_ok=True)


def extract_pptx(filepath):
    """提取 PPTX 文件中所有幻灯片的文字内容"""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    for idx, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_lines.append(text)
            # 尝试提取表格内容
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        ct = cell.text.strip()
                        if ct:
                            row_text.append(ct)
                    if row_text:
                        slide_lines.append(" | ".join(row_text))
        if slide_lines:
            slides_text.append(f"--- 第{idx}页 ---\n" + "\n".join(slide_lines))
    return "\n\n".join(slides_text)


def extract_pdf(filepath):
    """提取 PDF 文件文字内容"""
    from PyPDF2 import PdfReader
    reader = PdfReader(filepath)
    pages_text = []
    for idx, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages_text.append(f"--- 第{idx}页 ---\n" + text.strip())
    return "\n\n".join(pages_text)


def extract_docx(filepath):
    """提取 DOCX 文件文字内容"""
    from docx import Document
    doc = Document(filepath)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # 也提取表格
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                paragraphs.append(" | ".join(row_text))
    return "\n".join(paragraphs)


def extract_text(filepath, ext):
    """根据扩展名提取文件内容"""
    if ext in ('.txt', '.md'):
        with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    elif ext == '.pptx':
        return extract_pptx(filepath)
    elif ext == '.pdf':
        return extract_pdf(filepath)
    elif ext == '.docx':
        return extract_docx(filepath)
    else:
        return None


class ExamServer(http.server.SimpleHTTPRequestHandler):
    """服务器：静态文件 + 文件上传解析 API"""

    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=DIRECTORY, **kwargs)

    def end_headers(self):
        """添加禁用缓存的响应头"""
        self.send_header('Cache-Control', 'no-store, no-cache, must-revalidate, max-age=0')
        self.send_header('Pragma', 'no-cache')
        self.send_header('Expires', '0')
        super().end_headers()

    def log_message(self, format, *args):
        """只输出关键日志"""
        if args and 'POST' in str(args):
            print(f"  [API] {args[0]}")
        # 屏蔽普通 GET 请求日志

    def do_POST(self):
        if self.path == '/api/upload':
            try:
                # 解析 multipart/form-data
                content_type = self.headers.get('Content-Type', '')
                if 'multipart/form-data' not in content_type:
                    self.send_error(400, "需要 multipart/form-data")
                    return

                # 提取 boundary
                boundary = content_type.split('boundary=')[1].strip()
                if boundary.startswith('"') and boundary.endswith('"'):
                    boundary = boundary[1:-1]

                # 读取 body
                content_length = int(self.headers.get('Content-Length', 0))
                body = self.rfile.read(content_length)

                # 手动解析 multipart（避免 cgi.FieldStorage 的复杂性）
                filename, file_data = parse_multipart(body, boundary)

                if not filename or not file_data:
                    self.send_json(400, {"error": "未检测到上传文件"})
                    return

                # 保存临时文件
                ext = os.path.splitext(filename)[1].lower()
                tmp_path = os.path.join(UPLOAD_DIR, filename)
                with open(tmp_path, 'wb') as f:
                    f.write(file_data)

                # 提取内容
                try:
                    content = extract_text(tmp_path, ext)
                except Exception as e:
                    self.send_json(500, {"error": f"文件解析失败: {str(e)}"})
                    return
                finally:
                    # 清理临时文件
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)

                if content is None:
                    self.send_json(400, {"error": f"不支持的文件格式: {ext}"})
                    return

                if not content.strip():
                    self.send_json(400, {"error": "文件中未检测到文字内容，可能是纯图片型文件"})
                    return

                print(f"  [API] 成功解析 {filename}: {len(content)} 字")
                self.send_json(200, {
                    "filename": filename,
                    "content": content,
                    "wordCount": len(content),
                    "ext": ext
                })

            except Exception as e:
                print(f"  [API] 上传错误: {e}")
                self.send_json(500, {"error": f"服务器错误: {str(e)}"})
        else:
            self.send_error(404, "Not found")

    def do_OPTIONS(self):
        """CORS 预检"""
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def send_json(self, status, data):
        self.send_response(status)
        self.send_header('Content-Type', 'application/json; charset=utf-8')
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False).encode('utf-8'))


def parse_multipart(body, boundary):
    """手动解析 multipart/form-data"""
    boundary_bytes = boundary.encode('utf-8')
    parts = body.split(b'--' + boundary_bytes)

    for part in parts:
        if b'Content-Disposition' not in part:
            continue

        # 分离 header 和 body
        header_end = part.find(b'\r\n\r\n')
        if header_end == -1:
            continue

        header = part[:header_end].decode('utf-8', errors='replace')
        file_data = part[header_end + 4:]

        # 去掉末尾的 \r\n 和 boundary 尾部
        file_data = file_data.rstrip(b'\r\n')
        if file_data.endswith(b'--'):
            file_data = file_data[:-2].rstrip(b'\r\n')

        # 提取 filename
        filename_match = re.search(r'filename="(.+?)"', header)
        if filename_match:
            filename = filename_match.group(1)
            return filename, file_data

    return None, None


def find_free_port(preferred_port):
    """找可用端口"""
    try:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", preferred_port))
            return preferred_port
    except OSError:
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            s.bind(("127.0.0.1", 0))
            return s.getsockname()[1]


def open_browser(port):
    time.sleep(1)
    webbrowser.open(f"http://localhost:{port}")


if __name__ == "__main__":
    port = find_free_port(PORT)

    threading.Thread(target=open_browser, args=(port,), daemon=True).start()

    print()
    print("  ╔════════════════════════════════════════╗")
    print("  ║        智能出题系统 已启动（v2）       ║")
    print(f"  ║  地址: http://localhost:{port}          ║")
    print("  ║  支持格式: PPTX / PDF / DOCX / TXT     ║")
    print("  ║  关闭此窗口即可停止服务                ║")
    print("  ╚════════════════════════════════════════╝")
    print()

    try:
        with http.server.HTTPServer(("127.0.0.1", port), ExamServer) as httpd:
            httpd.serve_forever()
    except KeyboardInterrupt:
        print("\n服务已停止。")
        sys.exit(0)
